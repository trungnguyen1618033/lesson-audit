"""
assembler.py - Assemble captured slide images into a PPTX file.

Each PNG becomes one slide. Images are scaled to fill the slide while
maintaining aspect ratio and are centred on a 16:9 canvas.
"""

from pathlib import Path
from typing import List

from PIL import Image
from pptx import Presentation
from pptx.util import Emu, Inches, Pt


# Standard widescreen (16:9) dimensions in EMUs
SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)


def _emu_to_px(emu: int, dpi: int = 96) -> float:
    return emu / 914400 * dpi


def _fit_image(img_w: int, img_h: int, slide_w_emu: int, slide_h_emu: int):
    """
    Calculate left, top, width, height (EMU) to fit the image inside the slide
    while preserving aspect ratio and centering.
    """
    slide_ratio = slide_w_emu / slide_h_emu
    img_ratio = img_w / img_h

    if img_ratio >= slide_ratio:
        # Fit to width
        fit_w = slide_w_emu
        fit_h = int(slide_w_emu / img_ratio)
    else:
        # Fit to height
        fit_h = slide_h_emu
        fit_w = int(slide_h_emu * img_ratio)

    left = (slide_w_emu - fit_w) // 2
    top = (slide_h_emu - fit_h) // 2
    return left, top, fit_w, fit_h


def assemble_pptx(image_paths: List[Path], output_path: Path) -> Path:
    """
    Create a PPTX file from a list of image paths (in slide order).
    Returns the path to the created file.
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    blank_layout = prs.slide_layouts[6]  # completely blank layout

    for img_path in sorted(image_paths):
        slide = prs.slides.add_slide(blank_layout)

        with Image.open(img_path) as im:
            img_w, img_h = im.size

        left, top, width, height = _fit_image(
            img_w, img_h, int(SLIDE_WIDTH), int(SLIDE_HEIGHT)
        )
        slide.shapes.add_picture(str(img_path), left, top, width, height)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"PPTX saved: {output_path}")
    return output_path


def assemble_from_dir(image_dir: Path, name: str) -> Path:
    """
    Convenience wrapper: collect all slide_NNN.png files in `image_dir`
    and assemble them into <image_dir>/<name>.pptx.
    """
    images = sorted(image_dir.glob("slide_*.png"))
    if not images:
        raise FileNotFoundError(f"No slide images found in {image_dir}")
    output_path = image_dir / f"{name}.pptx"
    return assemble_pptx(images, output_path)
